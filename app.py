import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# =========================================================
# VALIDATION
# =========================================================
def validate_columns(df: pd.DataFrame):
    required = [
        "SKU",
        "Date",
        "Monthly_Demand",
        "Lead_Time_Days",
        "Shelf_Life_Days",
        "Unit_Cost",
        "On_Hand",
        "Expiration_Date",
        "Days_To_Expiry",
        "Expiring_Soon_Flag",
        "Inventory_Value",
    ]
    missing = [c for c in required if c not in df.columns]
    return missing


# =========================================================
# FORECAST + INVENTORY LOGIC (LIGHT MODEL)
# =========================================================
def compute_light_forecast(df: pd.DataFrame):
    df = df.copy()
    df["Date"] = pd.to_datetime(df["Date"])

    sku_list = df["SKU"].unique()
    results = []

    for sku in sku_list:
        sku_hist = df[df["SKU"] == sku].sort_values("Date")
        demand = sku_hist["Monthly_Demand"].astype(float)

        # Forecast: last 6-month avg * 12
        if len(demand) >= 6:
            avg_6 = demand.tail(6).mean()
        else:
            avg_6 = demand.mean()

        forecast_12m_units = float(avg_6 * 12)

        # Safety stock
        lead_time_days = float(sku_hist["Lead_Time_Days"].iloc[0])
        lead_time_months = max(lead_time_days / 30.0, 0.5)
        sigma = float(demand.std()) if demand.std() > 0 else 1.0
        Z = 1.65  # ~95% service level
        safety_stock_units = int(Z * sigma * np.sqrt(lead_time_months))

        # On-hand (use last non-null)
        sku_onhand = sku_hist[~sku_hist["On_Hand"].isna()]
        on_hand_units = float(sku_onhand["On_Hand"].iloc[-1]) if not sku_onhand.empty else 0.0

        # Cost
        unit_cost = float(sku_hist["Unit_Cost"].mean())

        # Expiry metrics (aggregate where On_Hand is present)
        exp_rows = sku_onhand.copy()
        if exp_rows.empty:
            total_on_hand = 0.0
            total_inv_value = 0.0
            exp_0_30_units = exp_31_90_units = exp_gt_90_units = 0.0
            exp_soon_units = exp_soon_value = 0.0
            min_days_to_expiry = np.nan
        else:
            total_on_hand = float(exp_rows["On_Hand"].sum())
            total_inv_value = float(exp_rows["Inventory_Value"].sum())
            min_days_to_expiry = float(exp_rows["Days_To_Expiry"].min())

            exp_0_30_units = float(exp_rows.loc[exp_rows["Days_To_Expiry"] <= 30, "On_Hand"].sum())
            exp_31_90_units = float(
                exp_rows.loc[
                    (exp_rows["Days_To_Expiry"] > 30) & (exp_rows["Days_To_Expiry"] <= 90),
                    "On_Hand",
                ].sum()
            )
            exp_gt_90_units = float(exp_rows.loc[exp_rows["Days_To_Expiry"] > 90, "On_Hand"].sum())

            exp_soon_units = exp_0_30_units + exp_31_90_units
            exp_soon_value = float(
                exp_rows.loc[exp_rows["Days_To_Expiry"] <= 90, "Inventory_Value"].sum()
            )

        # Recommended production
        recommended_prod_units = max(0, round(forecast_12m_units + safety_stock_units - on_hand_units))
        recommended_prod_value = recommended_prod_units * unit_cost

        # Risk label
        if total_on_hand == 0:
            risk_label = "No Stock"
        else:
            ratio = exp_soon_units / total_on_hand if total_on_hand > 0 else 0
            if ratio >= 0.5:
                risk_label = "High Expiry Risk"
            elif ratio >= 0.2:
                risk_label = "Moderate Expiry Risk"
            else:
                risk_label = "Low Expiry Risk"

        product_family = (
            sku_hist["Product_Family"].iloc[0]
            if "Product_Family" in sku_hist.columns
            else ""
        )

        results.append(
            {
                "SKU": sku,
                "Product_Family": product_family,
                "Lead_Time_Days": lead_time_days,
                "Shelf_Life_Days": float(sku_hist["Shelf_Life_Days"].iloc[0]),
                "Unit_Cost": unit_cost,
                # Forecast & stock
                "Forecast_12M_Units": round(forecast_12m_units),
                "Forecast_12M_Value": round(forecast_12m_units * unit_cost),
                "Safety_Stock_Units": safety_stock_units,
                "Safety_Stock_Value": round(safety_stock_units * unit_cost),
                "On_Hand_Units": round(on_hand_units),
                "On_Hand_Value": round(on_hand_units * unit_cost),
                # Expiry metrics
                "Total_On_Hand_Units": round(total_on_hand),
                "Total_On_Hand_Value": round(total_inv_value),
                "Expiring_0_30_Units": round(exp_0_30_units),
                "Expiring_31_90_Units": round(exp_31_90_units),
                "Expiring_>90_Units": round(exp_gt_90_units),
                "Expiring_Soon_Units_<=90": round(exp_soon_units),
                "Expiring_Soon_Value_<=90": round(exp_soon_value),
                "Min_Days_To_Expiry": min_days_to_expiry,
                # Recommendation
                "Recommended_Production_Units": recommended_prod_units,
                "Recommended_Production_Value": round(recommended_prod_value),
                # Risk label
                "Inventory_Risk_Label": risk_label,
            }
        )

    return pd.DataFrame(results)


# =========================================================
# EXCEL EXPORT
# =========================================================
def to_excel_bytes(df: pd.DataFrame) -> bytes:
    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="OptiStock_Output")
    return output.getvalue()


# =========================================================
# POWERPOINT EXPORT
# =========================================================
def build_powerpoint(results_df, insights, total_forecast_units,
                     total_recommended_value, total_at_risk_value, risk_pct):
    prs = Presentation()

    def add_header(slide, text):
        width = prs.slide_width
        rect = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=0,
            top=0,
            width=width,
            height=Inches(1.1),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(31, 42, 68)  # navy
        rect.line.color.rgb = RGBColor(31, 42, 68)
        tf = rect.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Slide 1 — Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "OptiStock – Executive Summary")

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8), Inches(2.5))
    tf = box.text_frame
    tf.text = "Inventory, Demand, and Expiry Analysis"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True

    # Slide 2 — KPIs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Key KPIs")

    bullets = [
        f"Total Forecasted Demand (12M): {total_forecast_units:,.0f} units",
        f"Total Recommended Production Value: ${total_recommended_value:,.0f}",
        f"Inventory Value at Risk (≤ 90 days): ${total_at_risk_value:,.0f}",
        f"% Inventory Value at Risk: {risk_pct:,.1f}%",
    ]

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.5), Inches(4))
    tf = body.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)

    # Slide 3 — Expiry Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Expiry Risk Overview")

    total_onhand = results_df["Total_On_Hand_Units"].sum()
    exp0 = results_df["Expiring_0_30_Units"].sum()
    exp1 = results_df["Expiring_31_90_Units"].sum()
    exp2 = results_df["Expiring_>90_Units"].sum()

    bullets = [
        f"Total On-Hand Units: {total_onhand:,.0f}",
        f"Units Expiring 0–30 Days: {exp0:,.0f}",
        f"Units Expiring 31–90 Days: {exp1:,.0f}",
        f"Units Expiring > 90 Days: {exp2:,.0f}",
        f"Value at Risk (≤ 90 Days): ${total_at_risk_value:,.0f}",
    ]

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.5), Inches(3.5))
    tf = body.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(16)

    # Slide 4 — Top 10 SKUs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Top 10 SKUs – Recommended Production")

    top10 = results_df.sort_values("Recommended_Production_Value", ascending=False).head(10)

    rows = len(top10) + 1
    table = slide.shapes.add_table(
        rows=rows,
        cols=4,
        left=Inches(0.5),
        top=Inches(1.5),
        width=Inches(9),
        height=Inches(3),
    ).table

    headers = ["SKU", "Units", "Value ($)", "Risk"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True

    for i, (_, row) in enumerate(top10.iterrows(), start=1):
        table.cell(i, 0).text = str(row["SKU"])
        table.cell(i, 1).text = f"{row['Recommended_Production_Units']:,.0f}"
        table.cell(i, 2).text = f"{row['Recommended_Production_Value']:,.0f}"
        table.cell(i, 3).text = row["Inventory_Risk_Label"]

    # Slide 5 — Key Insights
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Key Insights")

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.5), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True

    if insights:
        for i, line in enumerate(insights):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"- {line}"
            p.font.size = Pt(18)
    else:
        p = tf.paragraphs[0]
        p.text = "No major risks detected."
        p.font.size = Pt(18)

    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()


# =========================================================
# STREAMLIT APP LAYOUT
# =========================================================
st.set_page_config(page_title="OptiStock – Inventory Optimization", layout="wide")

# ---------- Global CSS for enterprise look ----------
st.markdown(
    """
    <style>
    :root {
        --primary-navy: #1F2A44;
        --primary-blue: #3E78B2;
        --light-gray: #F4F6FA;
        --card-bg: #FFFFFF;
        --border-gray: #C4CCD9;
        --risk-red: #D64545;
        --risk-yellow: #E3A008;
        --risk-green: #3BA55C;
    }

    /* Fix: Add more top padding so content doesn't overlap header */
    .block-container {
        padding-top: 2.5rem !important;
        padding-left: 2rem !important;
        padding-right: 2rem !important;
    }

    /* Top thin blue bar */
    .top-bar {
        height: 6px;
        width: 100%;
        background: var(--primary-blue);
    }

    /* Main navy header – reduced height so it no longer gets cut off */
    .main-header {
        background-color: var(--primary-navy);
        color: white;
        padding: 0.35rem 1.6rem 0.55rem 1.6rem; /* smaller top/bottom padding */
        box-shadow: 0 2px 4px rgba(0,0,0,0.15);
    }

    .main-header-title {
        font-size: 1.2rem; /* slightly smaller to fit on screen */
        font-weight: 700;
        margin-bottom: 0.1rem;
    }

    .main-header-subtitle {
        font-size: 0.82rem;
        margin-top: -0.1rem; /* pulls closer to title */
        opacity: 0.9;
    }

    /* Upload card styling */
    .upload-card {
        background-color: var(--card-bg);
        border-radius: 8px;
        border: 1px solid var(--border-gray);
        padding: 0.75rem 1rem;
        margin-top: 0.75rem;
        box-shadow: 0 1px 2px rgba(0,0,0,0.05);
    }

    /* Section headers */
    .section-title {
        font-size: 1.1rem;
        font-weight: 600;
        margin-top: 0.5rem;
        margin-bottom: 0.25rem;
        color: var(--primary-navy);
    }

    /* KPI Layout */
    .kpi-row {
        display: flex;
        gap: 0.75rem;
        flex-wrap: wrap;
        margin-top: 0.5rem;
        margin-bottom: 0.5rem;
    }

    .kpi-card {
        background-color: var(--card-bg);
        border-radius: 8px;
        border: 1px solid var(--border-gray);
        padding: 0.6rem 0.9rem;
        box-shadow: 0 1px 2px rgba(0,0,0,0.05);
        min-width: 180px;
        flex: 1 1 0;
    }

    .kpi-title {
        font-size: 0.8rem;
        text-transform: uppercase;
        letter-spacing: 0.06em;
        color: #6B7280;
        margin-bottom: 0.15rem;
    }

    .kpi-value {
        font-size: 1.2rem;
        font-weight: 700;
        color: #111827;
    }

    .kpi-sub {
        font-size: 0.8rem;
        color: #6B7280;
        margin-top: 0.1rem;
    }

    /* Buttons */
    .stDownloadButton button, .stButton button {
        border-radius: 999px;
        background-color: var(--primary-blue);
        color: white;
        border: 1px solid var(--primary-blue);
        padding: 0.35rem 0.9rem;
    }

    .stDownloadButton button:hover, .stButton button:hover {
        background-color: #325f8e;
        border-color: #325f8e;
        color: white;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

# ---------- Enterprise header ----------
st.markdown('<div class="top-bar"></div>', unsafe_allow_html=True)
st.markdown(
    """
    <div class="main-header">
        <div class="main-header-title">OptiStock™</div>
        <div class="main-header-subtitle">
            AI Inventory & Expiry Intelligence – Forecast, Safety Stock, and Risk in One View
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

# ---------- Upload section ----------
st.markdown('<div class="upload-card">', unsafe_allow_html=True)
st.markdown('<div class="section-title">Upload Combined Dataset</div>', unsafe_allow_html=True)
uploaded_file = st.file_uploader(
    "Upload your combined inventory & demand file (e.g., pharma_inventory_master_optistock.xlsx)",
    type=["xlsx"],
)
st.markdown("</div>", unsafe_allow_html=True)

if not uploaded_file:
    st.stop()

df = pd.read_excel(uploaded_file)

missing = validate_columns(df)
if missing:
    st.error(f"❌ Missing required columns: {missing}")
    st.stop()

df["Date"] = pd.to_datetime(df["Date"])
df["Expiration_Date"] = pd.to_datetime(df["Expiration_Date"])

results_df = compute_light_forecast(df)

# ---------------- EXEC SUMMARY METRICS & INSIGHTS ----------------
total_forecast_units = results_df["Forecast_12M_Units"].sum()
total_recommended_value = results_df["Recommended_Production_Value"].sum()
total_at_risk_value = results_df["Expiring_Soon_Value_<=90"].sum()
total_onhand_value = results_df["Total_On_Hand_Value"].sum()
risk_pct = (total_at_risk_value / total_onhand_value * 100) if total_onhand_value > 0 else 0.0

top_expiry = (
    results_df[results_df["Expiring_Soon_Value_<=90"] > 0]
    .sort_values("Expiring_Soon_Value_<=90", ascending=False)
    .head(5)
)
top_prod = results_df.sort_values("Recommended_Production_Value", ascending=False).head(5)

insights = []
if not top_expiry.empty:
    skus = ", ".join(top_expiry["SKU"].astype(str).tolist())
    share = (
        top_expiry["Expiring_Soon_Value_<=90"].sum() / total_at_risk_value * 100
        if total_at_risk_value > 0
        else 0
    )
    insights.append(
        f"Expiry risk (≤ 90 days) is concentrated in SKUs {skus}, which account for ~{share:,.0f}% of value at risk."
    )
if not top_prod.empty:
    skus = ", ".join(top_prod["SKU"].astype(str).tolist())
    insights.append(
        f"Recommended production value is primarily driven by SKUs {skus}, indicating key manufacturing focus areas."
    )
if total_onhand_value > 0:
    insights.append(
        f"Approximately {risk_pct:,.1f}% of total on-hand inventory value is at risk of expiry within 90 days."
    )
insights.append(
    f"Total 12-month forecasted demand across all SKUs is {total_forecast_units:,.0f} units."
)

# =========================================================
# TABS
# =========================================================
tab_summary, tab_forecast, tab_expiry, tab_top10, tab_download = st.tabs(
    ["📊 Executive Summary", "📈 Forecast & Plan", "⏳ Expiry Risk", "⭐ Top 10 SKUs", "📥 Download"]
)

# ---------------- EXEC SUMMARY TAB ----------------
with tab_summary:
    st.markdown('<div class="section-title">Executive Overview</div>', unsafe_allow_html=True)

    # KPI row
    kpi_html = f"""
    <div class="kpi-row">
      <div class="kpi-card">
        <div class="kpi-title">Forecasted Demand (12M)</div>
        <div class="kpi-value">{total_forecast_units:,.0f} units</div>
        <div class="kpi-sub">All SKUs combined</div>
      </div>
      <div class="kpi-card">
        <div class="kpi-title">Recommended Production Value</div>
        <div class="kpi-value">${total_recommended_value:,.0f}</div>
        <div class="kpi-sub">Next 12 months</div>
      </div>
      <div class="kpi-card">
        <div class="kpi-title">Value at Expiry Risk (≤ 90 days)</div>
        <div class="kpi-value">${total_at_risk_value:,.0f}</div>
        <div class="kpi-sub">Inventory at risk of write-off</div>
      </div>
      <div class="kpi-card">
        <div class="kpi-title">% Inventory Value at Risk</div>
        <div class="kpi-value">{risk_pct:,.1f}%</div>
        <div class="kpi-sub">Of total on-hand value</div>
      </div>
    </div>
    """
    st.markdown(kpi_html, unsafe_allow_html=True)

    st.markdown('<div class="section-title">Key Insights</div>', unsafe_allow_html=True)
    for i in insights:
        st.write(f"- {i}")

    st.markdown('<div class="section-title">Top Expiry Risk SKUs (by value ≤ 90 days)</div>', unsafe_allow_html=True)
    if not top_expiry.empty:
        st.dataframe(
            top_expiry[
                [
                    "SKU",
                    "Product_Family",
                    "Total_On_Hand_Units",
                    "Expiring_Soon_Units_<=90",
                    "Expiring_Soon_Value_<=90",
                    "Inventory_Risk_Label",
                ]
            ]
        )
    else:
        st.write("No SKUs with expiry risk ≤ 90 days in this dataset.")

# ---------------- FORECAST TAB ----------------
with tab_forecast:
    st.markdown('<div class="section-title">Forecast & Production Plan</div>', unsafe_allow_html=True)
    st.write(
        "Forecasted 12-month demand, safety stock, current on-hand inventory, "
        "and recommended production by SKU."
    )

    display_cols = [
        "SKU",
        "Product_Family",
        "Forecast_12M_Units",
        "Forecast_12M_Value",
        "Safety_Stock_Units",
        "Safety_Stock_Value",
        "On_Hand_Units",
        "On_Hand_Value",
        "Recommended_Production_Units",
        "Recommended_Production_Value",
        "Inventory_Risk_Label",
    ]
    st.dataframe(
        results_df[display_cols].sort_values("Recommended_Production_Value", ascending=False)
    )

# ---------------- EXPIRY TAB ----------------
with tab_expiry:
    st.markdown('<div class="section-title">Expiry Risk Dashboard</div>', unsafe_allow_html=True)

    total_on_hand_units = results_df["Total_On_Hand_Units"].sum()
    exp_0_30_units = results_df["Expiring_0_30_Units"].sum()
    exp_31_90_units = results_df["Expiring_31_90_Units"].sum()
    exp_gt_90_units = results_df["Expiring_>90_Units"].sum()

    c1, c2, c3 = st.columns(3)
    c1.metric("Total On-hand Units", f"{total_on_hand_units:,.0f}")
    c2.metric("Units Expiring ≤ 90 Days", f"{(exp_0_30_units + exp_31_90_units):,.0f}")
    c3.metric("Value Expiring ≤ 90 Days", f"${total_at_risk_value:,.0f}")

    st.markdown("#### Expiry Buckets (Units)")
    expiry_buckets = (
        pd.DataFrame(
            {
                "Bucket": ["0–30 days", "31–90 days", ">90 days"],
                "Units": [exp_0_30_units, exp_31_90_units, exp_gt_90_units],
            }
        )
        .set_index("Bucket")
    )
    st.bar_chart(expiry_buckets)

    st.markdown("#### SKUs with Highest Expiry Risk (Units ≤ 90 days)")
    top_expiry_units = (
        results_df[results_df["Expiring_Soon_Units_<=90"] > 0]
        .sort_values("Expiring_Soon_Units_<=90", ascending=False)
        .head(20)
    )
    st.dataframe(
        top_expiry_units[
            [
                "SKU",
                "Product_Family",
                "Total_On_Hand_Units",
                "Expiring_Soon_Units_<=90",
                "Expiring_Soon_Value_<=90",
                "Inventory_Risk_Label",
            ]
        ]
    )

# ---------------- TOP 10 TAB ----------------
with tab_top10:
    st.markdown('<div class="section-title">Top 10 SKUs by Recommended Production Value</div>', unsafe_allow_html=True)
    top10 = results_df.sort_values("Recommended_Production_Value", ascending=False).head(10)
    st.dataframe(
        top10[
            [
                "SKU",
                "Product_Family",
                "Recommended_Production_Units",
                "Recommended_Production_Value",
                "Total_On_Hand_Units",
                "Expiring_Soon_Units_<=90",
                "Expiring_Soon_Value_<=90",
                "Inventory_Risk_Label",
            ]
        ]
    )

# ---------------- DOWNLOAD TAB ----------------
with tab_download:
    st.markdown('<div class="section-title">Download Reports</div>', unsafe_allow_html=True)

    st.subheader("Excel Export")
    excel_bytes = to_excel_bytes(results_df)
    st.download_button(
        "📥 Download Excel",
        data=excel_bytes,
        file_name="optistock_results.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

    st.subheader("Executive PowerPoint Export")
    ppt_bytes = build_powerpoint(
        results_df,
        insights,
        total_forecast_units,
        total_recommended_value,
        total_at_risk_value,
        risk_pct,
    )
    st.download_button(
        "📊 Download Executive PowerPoint",
        data=ppt_bytes,
        file_name="optistock_executive_summary.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
